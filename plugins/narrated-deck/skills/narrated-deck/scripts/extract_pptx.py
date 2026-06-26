#!/usr/bin/env python3
"""extract_pptx.py — extract slide content from a PPTX into a structured outline.

Outputs a JSON document with one entry per slide:
  {
    "source": "deck.pptx",
    "slide_count": N,
    "slides": [
      {
        "idx": 1,
        "title": "Slide title",
        "subtitle": "Subtitle if any",
        "bullets": ["...", "..."],
        "notes": "Speaker notes (often the best source for narration)",
        "image_paths": ["extracted-images/slide-1-img-0.png", ...],
        "raw_text": "All text on the slide joined with newlines"
      }
    ]
  }

Use this output to drive Phase 2 (writing narration) — the speaker notes,
when present, are often a head-start on the script.

Dependencies:
  pip install python-pptx Pillow
"""

import argparse
import json
import sys
from pathlib import Path


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                  formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path", help="Path to .pptx file")
    ap.add_argument("--out", default="outline.json", help="Output JSON path")
    ap.add_argument("--images-dir", default="extracted-images",
                    help="Where to write extracted images (default: extracted-images/)")
    ap.add_argument("--skip-images", action="store_true",
                    help="Do not extract images (faster)")
    args = ap.parse_args()

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("ERROR: need python-pptx. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(2)

    src = Path(args.pptx_path)
    if not src.exists():
        print(f"ERROR: file not found: {src}", file=sys.stderr)
        sys.exit(2)

    prs = Presentation(str(src))
    img_dir = Path(args.images_dir)
    if not args.skip_images:
        img_dir.mkdir(parents=True, exist_ok=True)

    slides_out = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        subtitle = ""
        bullets = []
        raw_text_parts = []
        image_paths = []

        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                tf = shape.text_frame
                text = tf.text.strip()
                if not text:
                    continue
                raw_text_parts.append(text)

                # Heuristic: first text frame that's a placeholder labeled 'title' = title
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    if ph and ph.idx == 0 and not title:
                        title = text
                        continue
                    if ph and ph.idx == 1 and not subtitle:
                        subtitle = text
                        continue
                # Otherwise split paragraphs into bullets
                for para in tf.paragraphs:
                    pt = para.text.strip()
                    if pt and pt != title:
                        bullets.append(pt)

            # Images
            if not args.skip_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = img.ext or "png"
                    fname = img_dir / f"slide-{i:03d}-img-{len(image_paths)}.{ext}"
                    fname.write_bytes(img.blob)
                    image_paths.append(str(fname))
                except Exception as e:
                    print(f"  warn: could not extract image on slide {i}: {e}", file=sys.stderr)

        # Speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip() if notes_tf else ""

        # Fallback: if no title was found, use first non-empty bullet
        if not title and bullets:
            title = bullets.pop(0)

        slides_out.append({
            "idx": i,
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "notes": notes,
            "image_paths": image_paths,
            "raw_text": "\n".join(raw_text_parts),
        })

    out_path = Path(args.out)
    out_path.write_text(json.dumps({
        "source": str(src),
        "slide_count": len(slides_out),
        "slides": slides_out,
    }, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"✓ Extracted {len(slides_out)} slides → {out_path}")
    if not args.skip_images:
        n_imgs = sum(len(s["image_paths"]) for s in slides_out)
        print(f"  Images: {n_imgs} → {img_dir}/")
    n_notes = sum(1 for s in slides_out if s["notes"])
    print(f"  Slides with speaker notes: {n_notes}/{len(slides_out)}")
    if n_notes > 0:
        print("  ★ Speaker notes are usually the best starting point for narration.")


if __name__ == "__main__":
    main()
